import os
import tempfile
from datetime import date, datetime, timedelta
from decimal import Decimal

from langchain_core.tools import tool
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.db import SessionLocal
from app.model import Product, Bill, BillItem, BillStatus, Preference
from app.chart_utils import (
    sales_trend_chart,
    top_skus_chart,
    stock_health_chart,
    gst_collected_chart,
)
from app.tools.analytics_tools import _parse_date, _day_bounds_utc, _compute_live_day

DECK_OUTPUT_DIR = os.environ.get("DECK_OUTPUT_DIR", "/tmp/decks")

NAVY = RGBColor(0x2C, 0x3E, 0x50)
TEAL = RGBColor(0x16, 0xA0, 0x85)
GREY = RGBColor(0x7F, 0x8C, 0x8D)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _get_shop_name(db) -> str:
    prefs = db.query(Preference).filter(
        Preference.owner_id == "default",
        Preference.key == "shop_details",
    ).first()
    return (prefs.shop_name if prefs else None) or "Your Shop"


def _gather_range_data(db, start_date: date, end_date: date) -> dict:
    import json
    from app.model import DailyClosure

    dates: list[str] = []
    totals: list[float] = []
    range_sales = Decimal("0")
    range_tax = Decimal("0")
    range_bill_count = 0
    item_totals: dict[str, tuple[Decimal, Decimal]] = {}

    current = start_date
    while current <= end_date:
        closure = db.query(DailyClosure).filter(DailyClosure.closure_date == current).first()
        if closure:
            day_sales = Decimal(str(closure.total_sales))
            day_tax = Decimal(str(closure.total_tax))
            day_count = closure.bill_count
            day_items = json.loads(closure.top_items_json or "[]")
        else:
            day_data = _compute_live_day(db, current)
            day_sales = day_data["total_sales"]
            day_tax = day_data["total_tax"]
            day_count = day_data["bill_count"]
            day_items = day_data["top_items"]

        dates.append(current.strftime("%d %b"))
        totals.append(float(day_sales))
        range_sales += day_sales
        range_tax += day_tax
        range_bill_count += day_count
        for name, qty, revenue in day_items:
            q, r = item_totals.get(name, (Decimal("0"), Decimal("0")))
            item_totals[name] = (q + Decimal(str(qty)), r + Decimal(str(revenue)))

        current += timedelta(days=1)

    top_items = [
        (name, qty, revenue)
        for name, (qty, revenue) in sorted(item_totals.items(), key=lambda kv: kv[1][1], reverse=True)[:8]
    ]

    return {
        "dates": dates,
        "totals": totals,
        "range_sales": range_sales,
        "range_tax": range_tax,
        "range_bill_count": range_bill_count,
        "top_items": top_items,  # [(name, qty, revenue), ...]
    }


def _gather_gst_by_slab(db, start_date: date, end_date: date) -> dict:
    start_utc, _ = _day_bounds_utc(start_date)
    _, end_utc = _day_bounds_utc(end_date)

    rows = (
        db.query(BillItem.gst_slab, BillItem.cgst_amt, BillItem.sgst_amt)
        .join(Bill, BillItem.bill_id == Bill.id)
        .filter(
            Bill.status == BillStatus.finalized,
            Bill.finalized_at >= start_utc,
            Bill.finalized_at < end_utc,
        )
        .all()
    )
    by_slab: dict[float, Decimal] = {}
    for slab, cgst, sgst in rows:
        slab = float(slab)
        by_slab[slab] = by_slab.get(slab, Decimal("0")) + Decimal(str(cgst)) + Decimal(str(sgst))

    slabs = sorted(by_slab.keys())
    return {
        "labels": [f"{s:g}%" for s in slabs],
        "amounts": [float(by_slab[s]) for s in slabs],
    }


def _gather_stock_health(db, limit: int = 10) -> dict:
    products = db.query(Product).all()
    if not products:
        return {"names": [], "qty": [], "thresholds": [], "low_stock_count": 0}

    ranked = sorted(products, key=lambda p: float(p.qty_on_hand) - float(p.reorder_level))
    low_stock_count = sum(1 for p in products if float(p.qty_on_hand) <= float(p.reorder_level))
    worst = ranked[:limit]

    return {
        "names": [p.name for p in worst],
        "qty": [float(p.qty_on_hand) for p in worst],
        "thresholds": [float(p.reorder_level) for p in worst],
        "low_stock_count": low_stock_count,
    }


def _add_title_slide(prs: Presentation, start_date: date, end_date: date, shop_name: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.text = "Sales Analysis"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.7))
    tf2 = subtitle_box.text_frame
    tf2.text = f"{shop_name} — {start_date.strftime('%d %b %Y')} to {end_date.strftime('%d %b %Y')}"
    tf2.paragraphs[0].font.size = Pt(20)
    tf2.paragraphs[0].font.color.rgb = TEAL

    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
    tf3 = footer_box.text_frame
    tf3.text = f"Generated {datetime.utcnow().strftime('%d %b %Y')}"
    tf3.paragraphs[0].font.size = Pt(12)
    tf3.paragraphs[0].font.color.rgb = GREY


def _add_chart_slide(prs: Presentation, title: str, image_path: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY

    # chart_utils saves at figsize=(9, 4.5) => keep aspect ratio, center it
    pic_width = Inches(10.5)
    left = (SLIDE_W - pic_width) / 2
    slide.shapes.add_picture(image_path, left, Inches(1.3), width=pic_width)


def _add_insights_slide(prs: Presentation, data: dict, gst_data: dict, stock_data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Key Insights"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY

    range_sales = data["range_sales"]
    range_tax = data["range_tax"]
    bill_count = data["range_bill_count"]
    avg_bill = (range_sales / bill_count) if bill_count else Decimal("0")
    top_item_name = data["top_items"][0][0] if data["top_items"] else "—"
    gst_total = sum(gst_data["amounts"]) if gst_data["amounts"] else 0.0

    bullets = [
        f"Total sales: ₹{range_sales:,.2f} across {bill_count} bill(s)",
        f"Tax collected: ₹{range_tax:,.2f} (GST breakup: {', '.join(f'{l} ₹{a:,.0f}' for l, a in zip(gst_data['labels'], gst_data['amounts'])) or 'none'})",
        f"Average bill value: ₹{avg_bill:,.2f}",
        f"Top-selling item: {top_item_name}",
        f"Products at/below reorder level: {stock_data['low_stock_count']}",
    ]

    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5))
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
        p.text = f"•  {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = NAVY
        p.space_after = Pt(14)


@tool
def generate_report_pptx(start: str, end: str) -> str:
    """Generate a business-analysis PowerPoint deck for a date range: a title
    slide, sales trend chart, top-selling items chart, stock health chart,
    GST-collected-by-slab chart, and a key-insights summary slide. Dates are
    format YYYY-MM-DD, inclusive. Use this whenever the owner asks for a
    sales analysis deck, a weekly/monthly report, or similar — e.g. "make
    this week's sales analysis deck". Returns the file path of the generated
    .pptx so it can be sent to the owner (e.g. via Telegram sendDocument).
    Reads the same closed/live day data as get_sales_range, so the numbers
    in the deck always match what get_daily_summary and get_sales_range
    would report for the same range."""
    db = SessionLocal()
    tmp_files: list[str] = []
    try:
        start_date = _parse_date(start)
        end_date = _parse_date(end)
        if not start_date or not end_date:
            return "Couldn't parse one of the dates — use YYYY-MM-DD for both start and end."
        if start_date > end_date:
            return f"start ({start}) is after end ({end}) — did you mean to swap them?"
        if (end_date - start_date).days > 366:
            return "Range is over a year — narrow it down for a meaningful deck."

        range_data = _gather_range_data(db, start_date, end_date)
        gst_data = _gather_gst_by_slab(db, start_date, end_date)
        stock_data = _gather_stock_health(db)

        if range_data["range_bill_count"] == 0:
            return f"No finalized sales between {start} and {end} — nothing to build a deck from yet."

        shop_name = _get_shop_name(db)

        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        _add_title_slide(prs, start_date, end_date, shop_name)

        with tempfile.TemporaryDirectory() as tmp_dir:
            # Sales trend
            trend_path = os.path.join(tmp_dir, "trend.png")
            sales_trend_chart(range_data["dates"], range_data["totals"], trend_path)
            _add_chart_slide(prs, "Daily Sales Trend", trend_path)

            # Top-selling items (by revenue)
            if range_data["top_items"]:
                names = [name for name, qty, revenue in range_data["top_items"]]
                revenues = [float(revenue) for name, qty, revenue in range_data["top_items"]]
                top_path = os.path.join(tmp_dir, "top_skus.png")
                top_skus_chart(names, revenues, top_path)
                _add_chart_slide(prs, "Top-Selling Items", top_path)

            # Stock health
            if stock_data["names"]:
                stock_path = os.path.join(tmp_dir, "stock_health.png")
                stock_health_chart(stock_data["names"], stock_data["qty"], stock_data["thresholds"], stock_path)
                _add_chart_slide(prs, "Stock Health", stock_path)

            # GST collected by slab
            if gst_data["labels"]:
                gst_path = os.path.join(tmp_dir, "gst.png")
                gst_collected_chart(gst_data["labels"], gst_data["amounts"], gst_path)
                _add_chart_slide(prs, "GST Collected by Slab", gst_path)

            _add_insights_slide(prs, range_data, gst_data, stock_data)

            os.makedirs(DECK_OUTPUT_DIR, exist_ok=True)
            output_path = os.path.join(
                DECK_OUTPUT_DIR, f"analysis_{start_date.isoformat()}_to_{end_date.isoformat()}.pptx"
            )
            prs.save(output_path)

        return f"FILE_PATH: {output_path}"
    finally:
        db.close()